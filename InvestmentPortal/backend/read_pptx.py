# -*- coding: utf-8 -*-
"""Read and print PPTX content"""
from pptx import Presentation
import sys
sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation(r'D:\Industry\산업자료\6. 전력 인프라\전력 인프라.pptx')
print(f"Total slides: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    print(f'\n=== Slide {i+1} ===')
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(repr(shape.text[:800]))
    print()
