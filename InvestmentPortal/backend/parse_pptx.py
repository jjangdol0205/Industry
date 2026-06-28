from pptx import Presentation
import os

path = r'D:\Industry\산업자료\6. 전력 인프라\전력 인프라.pptx'
prs = Presentation(path)

print(f"Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides):
    print(f"=== Slide {i+1} ===")
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            text = shape.text.strip()
            if len(text) > 20:
                print(text[:600])
    print()
